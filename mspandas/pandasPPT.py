import pptx
import pandas as pd
import numpy as np

from mspandas.utils.ppt import pptFunctions
from mspandas.utils.pd import pdFunctions


class Table():
    """
    Abstract table class to convert a Pandas DataFrame into a PowerPoint table.

    Attributes
    ----------
    shape: pptx.shapes.placeholder.TablePlaceholder
        empty pptx table shape, placeholder object for table graphic frame.
    df: pd.DataFrame
        Pandas DataFrame to be converted into table.

    Parameters
    ----------
    header: bool
        Whether or not to include DataFrame's header in table representation.
    index: bool
        Whether or not to include DataFrame's index in table representation.
    keep_names: str
        Priority for keeping names on axis when transforming. If table includes header and index, and DataFrame names specified for both axis, will keep names based on priority.
        Specificy 'index' for column-wise totals, 1 or 'columns' for row-wise totals.
    na_rep: int,float,str
        String value for representing null values.
    dtype_format: dict
        Map of numpy data types to string format.
    font_size: int
        Cell text font size. For more, see pptx.util.Pt
    font_color: tuple
        Cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    font_name:
        Cell text font name, for example 'Arial'.
    column_totals: bool
        Whether or not to include column totals in table representation.
    column_totals_label: str
        Index label for totals row, can be string or tuple for multiindex.
    column_totals_method: str/func
        Function to use for aggregating the columns.
    column_totals_aggmap: dict
        Map of column names and aggregation method to be applied when totaling columns. Example: {'a':'mean', 'b':'mean'}. Default will sum all data.
	row_totals: bool
        Whether or not to include row totals in table representation.
    row_totals_label: str
        Column label for totals column, can be string or tuple for multiindex header.
    row_totals_method: str/func
        Function to use for aggregating the rows.
    row_totals_aggmap: dict
        Map of index names and aggregation method to be applied when totaling columns. Example: {0:'mean', 1:'mean'}. Default will sum all data.
    fill_header: bool
        Whether or not to fill the cell backgound color of table header rows.
    bold_header: bool
        Whether or not to bold the text in all table header rows.
    header_font_size: int
        Header cell text font size. For more, see pptx.util.Pt
    header_font_color: tuple
        Header cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    fill_index: bool
        Whether or not to fill the cell backgound color of table index columns.
    bold_index: bool
        Whether or not to bold the text in all table index columns.
    index_font_size: int
        Index cell text font size. For more, see pptx.util.Pt
    index_font_color: tuple
        Index cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
        Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
    row_banding: bool
        Whether or not turn on default PowerPoint horizontal banding styles.
    column_banding: bool
        Whether or not turn on default PowerPoint vertical banding styles.
    first_row: bool
        Whether of not to turn on default PowerPoint styles for first row.
    first_col: bool
        Whether of not to turn on default PowerPoint styles for first col.
    last_row: bool
        Whether of not to turn on default PowerPoint styles for last row.
    last_row: bool
        Whether of not to turn on default PowerPoint styles for last row.
    merge_indices: bool
        Whether or not to merge adjacent equal values in table indices (rows in index columns or columns in header rows).
    center_merge: str
        Whether or not to center the paragraph text after merge.
    cell_margins: str
		Keyword for setting cell margin widths. Use one of 'normal', 'none', 'narrow', 'tight', or 'wide'. Keywords are adopted from ppt with a custom tight setting.
    row_height: float
        Row height in inches.
    align_num_cols: str
        Paragraph alignmnet applied as default to all numeric columns. Default 'center'.
    align_char_cols: str
        Paragraph alignmnet applied as default to all string columns. Default 'left'.
    align_dt_cols: str
        Paragraph alignmnet applied as default to all datetime columns. Default 'center'.
    align_col_map: dict
        Paragraph alignment applied on custom columns by DataFrame column name. Example: {'Column One':'left', 'Column Two':'center'}. For me see pptx.enum.text.PP_ALIGN

    Methods
    -------
    add_totals: Aggregate data in DataFrame, applied as column-wise or row-wise by axis argument.
    transform: Transofrm DataFrame to mirror output representation.
    format_index: Format DataFrame index values as unicode strings.
    format_values: Format DataFrame values as unicode strings.
    insert_table: Insert a graphic frame table object into the table placeholder object.
    style_index: Apply styles to DataFrame index in PowerPoint table in place.
    align_columns: Apply column alignment to table in place.
    style_table: Apply styles to PowerPoint table in place.
    convert: Perform the conversion from DataFrame values to table cells.
    """

    def __init__(self, shape, df, **kwargs):
        self.shape = shape
        self.df = df

        self.header = kwargs.pop('header',False)
        self.index = kwargs.pop('index',False)
        self.keep_names = kwargs.pop('keep_names','columns')

        self.na_rep = kwargs.pop('na_rep',' ')
        self.dtype_format = kwargs.pop('dtype_format',None)

        self.font_size = kwargs.pop('font_size',None)
        self.font_color = kwargs.pop('font_color',None)
        self.font_name = kwargs.pop('font_name',None)

        self.column_totals = kwargs.pop('column_totals',False)
        self.column_totals_label = kwargs.pop('column_totals_label','Total')
        self.column_totals_method = kwargs.pop('column_totals_method',np.sum)
        self.column_totals_aggmap = kwargs.pop('column_totals_aggmap',{})
        self.row_totals = kwargs.pop('row_totals',False)
        self.row_totals_label = kwargs.pop('row_totals_label','Total')
        self.row_totals_method = kwargs.pop('row_totals_method',np.sum)
        self.row_totals_aggmap = kwargs.pop('row_totals_aggmap',{})

        self.fill_header = kwargs.pop('fill_header',True)
        self.bold_header = kwargs.pop('bold_header',True)
        self.header_font_size = kwargs.pop('header_font_size',None)
        self.header_font_color = kwargs.pop('header_font_color',None)
        self.fill_index = kwargs.pop('fill_index',False)
        self.bold_index = kwargs.pop('bold_index',True)
        self.index_font_size = kwargs.pop('index_font_size',None)
        self.index_font_color = kwargs.pop('index_font_color',None)
        self.fill_color = kwargs.pop('fill_color',pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1)

        self.row_banding = kwargs.pop('row_banding',True)
        self.column_banding = kwargs.pop('column_banding',False)
        self.first_row = kwargs.pop('first_row',False)
        self.first_col = kwargs.pop('first_col',False)
        self.last_row = kwargs.pop('last_row',False)
        self.last_col = kwargs.pop('last_col',False)

        self.merge_indices = kwargs.pop('merge_indices',True)
        self.center_merge = kwargs.pop('center_merge',True)

        self.cell_margins = kwargs.pop('cell_margins','tight')

        self.row_height = kwargs.pop('row_height',0.15)

        self.align_num_cols = kwargs.pop('align_num_cols','center')
        self.align_char_cols = kwargs.pop('align_char_cols','left')
        self.align_dt_cols = kwargs.pop('align_dt_cols','center')
        self.align_col_map = kwargs.pop('align_col_map',{})

    def add_totals(self, **kwargs):
        """Aggregate data in DataFrame, applied as column-wise or row-wise by axis argument.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to have columns totaled.
        totals_label: str/tuple
            Index (or columns) label for totals row (or column), can be string or tuple for multiindex.
        totals_method: str/func
            Function to use for aggregating data. Default will apply np.sum.
        totals_aggmap: dict
            Map of column (or index) names and aggregation method to be applied when totaling. Example: {'a':'mean', 'b':'mean'}. Takes priority over total_method.
        axis: int
            Orientation, based on axis of DataFrame index, for totaling. 0 or 'index' for column-wise totals, 1 or 'columns' for row-wise totals.

        Returns
        -------
        data: pandas.DataFrame
            DataFrame with new row as column totals.
        """
        data = kwargs.pop('data',self.df.copy())
        totals_label = kwargs.pop('totals_label','Totals')
        totals_method = kwargs.pop('totals_method',np.sum)
        totals_aggmap = kwargs.pop('totals_aggmap',{})
        axis = kwargs.pop('axis',0)
        if axis in [1,'columns']:
            data = data.T
        elif axis in [0,'index']:
            pass
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for column-wise totals, 1 or "columns" for row-wise totals.'))
        names = list(data.index.names)
        ordered_columns = list(data.columns)
        # totals
        c_totals = data.fillna(0).agg({col:totals_aggmap.get(col,totals_method) for col in data.columns}).rename(totals_label)
        c_totals = c_totals.to_frame().T
        # create multiindex if needed
        for i in range(data.index.nlevels-1):
            c_totals['dummy'] = ' '
            c_totals = c_totals.set_index('dummy',append=True)
            c_totals.index.names = [None]*len(c_totals.index.names)
        try:
        	data = pd.concat([data, c_totals], axis=0)
        except TypeError:
        	# df index is categorical
        	# add label as category and append
        	data.index = data.index.add_categories(totals_label)
        	data = pd.concat([data, c_totals], axis=0)
        data = data.reindex(ordered_columns, axis=1)
        data.index.names = names
        if axis in [1,'columns']:
            return data.T
        else:
            return data

    def transform(self, **kwargs):
        """Transofrm DataFrame into output representation with presentation options.

        Optional header, index or totals are converted into DataFrame row values.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        header: bool
            Whether or not to include DataFrame's header in table representation.
        index: bool
            Whether or not to include DataFrame's index in table representation.
        keep_names: str
            Priority for keeping names on axis when transforming. If table includes header and index, and DataFrame names specified for both axis, will keep names based on priority.
            Specificy 'index' for column-wise totals, 1 or 'columns' for row-wise totals.

        Returns
        -------
        data: pandas.DataFrame
            Transformed DataFrame including optional header, index and totals.

        Notes:
        ------
        When header is True and DataFrame column names are strings, \
        Series of non-object dtypes will be converted to having dtypes of object.
        """
        data = kwargs.pop('data',self.df.copy())
        header = kwargs.pop('header',self.header)
        index = kwargs.pop('index',self.index)
        keep_names = kwargs.pop('keep_names',self.keep_names)
        if keep_names == 'index':
            if index:
                data = data.reset_index()
            if header:
                data = data.T.reset_index().T
        elif keep_names == 'columns':
            if header:
                data = data.T.reset_index().T
            if index:
                data = data.reset_index()
        else:
            raise(ValueError('Incorrect value for keep_names. Use "index", or "columns"'))
        return data

    def format_index(self, **kwargs):
        """Format DataFrame index values as unicode strings.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        dtype_format: dict
            Map of numpy data types to string format.
        axis: int
            Axis of DataFrame index to be formatted. 0 or 'index' for index, 1 or 'columns' for header.

        Returns
        -------
        data: pandas.DataFrame
            Formatted DataFrame where all index values are dtype np.unicode.
        """
        data = kwargs.pop('data',self.df.copy())
        dtype_format = kwargs.pop('dtype_format',self.dtype_format)
        axis = kwargs.pop('axis',0)
        if axis in [1,'columns']:
            index = data.columns
        elif axis in [0,'index']:
            index = data.index
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for index, 1 or "columns" for header.'))
        names=index.names
        index_vals=[]
        for n in range(index.nlevels):
            vals = index.get_level_values(n)
            if dtype_format is not None:
                for dtype,fmt in dtype_format.items():
                    if np.issubdtype(vals.dtype,dtype):
                        if np.issubdtype(dtype,np.datetime64):
                            vals = vals.strftime(fmt)
                        elif np.issubdtype(dtype,np.number):
                            vals = vals.format(fmt)
                        else:
                            raise NotImplementedError('Not able to convert values of dtype {} to strings.\
                            Convert manually in your DataFrame before passing into Table()'.format(dtype))
            vals = vals.astype(np.unicode)
            index_vals.append(vals)
        index = pd.MultiIndex.from_arrays(index_vals,names=names)
        if axis in [1,'columns']:
            data.columns = index
        elif axis in [0,'index']:
            data.index = index
        return data

    def format_values(self, **kwargs):
        """Format DataFrame values as unicode strings.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        dtype_format: dict
            Map of numpy data types to string format.

        Returns
        -------
        data: pandas.DataFrame
            Formatted DataFrame where all values are dtype np.unicode.
        """
        data = kwargs.pop('data',self.df.copy())
        dtype_format = kwargs.pop('dtype_format',self.dtype_format)
        if dtype_format is not None:
            for dtype,fmt in dtype_format.items():
                for col,x in data.iteritems():
                    if np.issubdtype(x.dtype,dtype):
                        if np.issubdtype(dtype,np.datetime64):
                            data.loc[:,col] = x.dt.strftime(fmt)
                        elif np.issubdtype(dtype,np.number):
                            data.loc[:,col] = x.apply(fmt.format)
                        else:
                            raise NotImplementedError('Not able to convert values of dtype {} to strings.\
                            Convert manually in your DataFrame before passing into Table()'.format(dtype))
        data = data.fillna(self.na_rep)
        data = data.astype(np.unicode)
        return data

    def insert_table(self):
        """Insert a graphic frame table object into the table placeholder.

        Notes
        -----
        Placeholders become invalid after inserting a graphic frame object. The new shape is the return value of the insert_table() call and may also be obtained from the placeholders collection using the same idx key.
        For more on this see https://python-pptx.readthedocs.io/en/latest/user/placeholders-using.html#insert-content-into-a-placeholder

        """
        rows,cols = self.transform().shape
        if self.column_totals:
            rows+=1
        if self.row_totals:
            cols+=1
        if self.shape.is_placeholder and (not self.shape.has_table and not self.shape.shape_id == None):
            self.shape = self.shape.insert_table(rows=rows,
                                                 cols=cols)
        else:
            raise Warning('Shape object is not a placeholder or already contains a table graphic frame.')

    def style_index(self, **kwargs):
        """Apply styles to DataFrame index in PowerPoint table in place.

        Parameters
        ----------
        font_size: int
            Cell text font size. For more, see pptx.util.Pt
        font_color: tuple
            Cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
        bold: bool
            Whether or not to bold the text.
        fill: bool
            Whether or not to fill the cell backgound color.
        fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
            Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
        merge_indices: bool
            Whether or not to merge adjacent equal values in table indices (rows in index columns or columns in header rows).
        center_merge: str
            Whether or not to center the paragraph text after merge.
        axis: int
            Axis of DataFrame index to be formatted. 0 or 'index' for index, 1 or 'columns' for header.

        """
        table = self.shape.table
        font_size = kwargs.pop('font_size',None)
        font_color = kwargs.pop('font_color',None)
        bold = kwargs.pop('bold',False)
        fill = kwargs.pop('fill',True)
        fill_color = kwargs.pop('fill_color',self.fill_color)
        merge_indices = kwargs.pop('merge_indices',self.merge_indices)
        center_merge = kwargs.pop('center_merge',self.center_merge)
        axis = kwargs.pop('axis',0)
        data = self.df.copy()
        rows,cols = self.transform().shape
        if axis in [1,'columns']:
            index = data.columns
            axis = 1
            offset = cols - len(index)
            numcells = cols+1 if self.row_totals else cols
        elif axis in [0,'index']:
            index = data.index
            axis = 0
            offset = rows - len(index)
            numcells = rows+1 if self.column_totals else rows
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for index, 1 or "columns" for header.'))
        for n in range(index.nlevels):
            merge = False
            for i in range(numcells):
                if axis==0:
                    c = table.cell(i,n)
                else:
                    c = table.cell(n,i)
                ### Start Section on Auto-Merging (unstable feature)
                if merge_indices and i >= offset:
                    # table loc to dataframe loc
                    j = i-offset
                    if j < len(index)-1:
                        equal = True if index.get_level_values(n)[j] == index.get_level_values(n)[j+1] else False
                        if not merge and equal:
                            merge = True
                            mergestart = i
                            # default to maximum merge span
                            mergespan = len(index)-j
                            # calc actual mergespan using next non-equal value
                            if len(index) > j+1:
                                for k,d in enumerate(index.get_level_values(n)[j+1:]):
                                    if not index.get_level_values(n)[j] == d:
                                        mergespan = k+1
                                        break
                            if axis in [0,'rows']:
                                c._tc.set('rowSpan', str(mergespan))
                            else:
                                c._tc.set('gridSpan', str(mergespan))
                            if center_merge:
                                tf = c.text_frame
                                p = tf.paragraphs[0]
                                p.alignment = pptx.enum.text.PP_ALIGN.__dict__['CENTER']
                        elif merge and not equal:
                            # stop merge
                            if axis in [0,'rows']:
                                c._tc.set('vMerge', '1')
                            else:
                                c._tc.set('hMerge', '1')
                            merge = False
                    elif merge and i > mergestart and i < mergestart+mergespan-1:
                        # continue merge
                        if axis in [0,'rows']:
                            c._tc.set('vMerge', '1')
                        else:
                            c._tc.set('hMerge', '1')
                    elif merge and i > mergestart and i == mergestart+mergespan-1:
                        # end of merge
                        if axis in [0,'rows']:
                            c._tc.set('vMerge', '1')
                        else:
                            c._tc.set('hMerge', '1')
                        merge = False
                ### End Section on Auto-Merging
                if fill:
                    c.fill.solid()
                    if isinstance(fill_color,pptx.enum.base.EnumValue):
                        c.fill.fore_color.theme_color = fill_color
                    elif isinstance(fill_color,tuple):
                        c.fill.fore_color.rgb = pptx.dml.color.RGBColor(*fill_color)
                    elif isinstance(fill_color,str):
                        c.fill.fore_color.rgb = pptx.dml.color.RGBColor.from_string(fill_color) if not fill_color.startswith('#') else pptx.dml.color.RGBColor.from_string(fill_color[1:])
                    else:
                        raise ValueError('Incorrect value for fill_color. \
                        Please provide one of RGB code as `tuple` of 3 integers, HEX code as string, or an instance of `pptx.enum.dml.MSO_THEME_COLOR`')
                c = pptFunctions.format_cell(c,
                                          fill=fill,
                                          fill_color=fill_color,
                                          font_size=font_size,
                                          font_color=font_color,
                                          bold=bold)

    def align_columns(self, **kwargs):
        """Apply column alignment to table in place.

        Parameters
        ----------
        align_num_cols: str
            Paragraph alignmnet applied as default to all numeric columns. Default 'center'.
        align_char_cols: str
            Paragraph alignmnet applied as default to all string columns. Default 'left'.
        align_dt_cols: str
            Paragraph alignmnet applied as default to all datetime columns. Default 'center'.
        align_col_map: dict
            Paragraph alignment applied on custom columns by DataFrame column name. Example: {'Column One':'left', 'Column Two':'center'}. For me see pptx.enum.text.PP_ALIGN

        """
        align_num_cols = kwargs.pop('align_num_cols',self.align_num_cols)
        align_char_cols = kwargs.pop('align_char_cols',self.align_char_cols)
        align_dt_cols = kwargs.pop('align_dt_cols',self.align_dt_cols)
        align_col_map = kwargs.pop('align_col_map',self.align_col_map)
        table = self.shape.table
        data = self.df.copy()
        rows,cols = self.transform().shape
        col_offset = cols - len(data.columns)
        row_offset = rows - len(data.index)
        num_col_cells = cols+1 if self.row_totals else cols
        num_row_cells = rows+1 if self.column_totals else rows
        for j in range(num_col_cells):
            if j < col_offset:
                # index columns
                alignment = 'left'
            elif j >= len(data.columns):
                # totals column
                alignment = 'center'
            else:
                name = data.columns[j]
                dtype = data.loc[:,name].values.dtype
                if name in align_col_map.keys():
                    alignment = align_col_map[name]
                else:
                    if np.issubdtype(dtype,np.number):
                        alignment =  align_num_cols
                    elif np.issubdtype(dtype,np.char):
                        alignment = align_char_cols
                    elif np.issubdtype(dtype,np.datetime64):
                        alignment = align_dt_cols
                    else:
                        raise Warning('No default alignment defined for columns to dtype {}.\
                        Alignment will be PowerPoint default'.format(dtype))
            for i in range(num_row_cells):
                if i < row_offset:
                    # skip header rows
                    continue
                c = table.cell(i,j)
                tf = c.text_frame
                p = tf.paragraphs[0]
                if isinstance(alignment,pptx.enum.text.PP_ALIGN):
                    p.alignment = alignment
                elif isinstance(alignment,str):
                    p.alignment = pptx.enum.text.PP_ALIGN.__dict__[alignment.upper()]
                else:
                    raise ValueError("Incorrect value for alignment. \
                    Please provide a string like 'center' or 'left', or an instance of `pptx.enum.text.PP_ALIGN`")

    def style_table(self, **kwargs):
        """Apply styles to PowerPoint table in place.

        Parameters
        ----------
        header: bool
            Whether or not to include DataFrame's header in table representation.
        fill_header: bool
            Whether or not to fill the cell backgound color of table header rows.
        bold_header: bool
            Whether or not to bold the text in all table header rows.
        header_font_size: int
            Header cell text font size. For more, see pptx.util.Pt
        header_font_color: tuple
            Header cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
        index: bool
            Whether or not to include DataFrame's index in table representation.
        fill_index: bool
            Whether or not to fill the cell backgound color of table index columns.
        bold_index: bool
            Whether or not to bold the text in all table index columns.
        index_font_size: int
            Index cell text font size. For more, see pptx.util.Pt
        index_font_color: tuple
            Index cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
        fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
            Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
        row_banding: bool
            Whether or not turn on default PowerPoint horizontal banding styles.
        column_banding: bool
            Whether or not turn on default PowerPoint vertical banding styles.
        first_row: bool
            Whether of not to turn on default PowerPoint styles for first row.
        first_col: bool
            Whether of not to turn on default PowerPoint styles for first col.
        last_row: bool
            Whether of not to turn on default PowerPoint styles for last row.
        last_row: bool
            Whether of not to turn on default PowerPoint styles for last row.
        merge_indices: bool
            Whether or not to merge adjacent equal values in table indices (rows in index columns or columns in header rows).
        center_merge: str
            Whether or not to center the paragraph text after merge.
        row_height: float
			Row height in inches.

        Notes
        -----
        We do not apply the PowerPoint styles, controlling overarching table theme, as they are not currently supported by python-pptx (see https://github.com/scanny/python-pptx/issues/27)
        Instead we apply the logical DataFrame styling, emphasizing header, index and totals with bold text or filled backgound
        """
        table =  self.shape.table
        fill_header = kwargs.pop('fill_header',self.fill_header)
        bold_header = kwargs.pop('bold_header',self.bold_header)
        header_font_size = kwargs.pop('header_font_size',self.header_font_size)
        header_font_color = kwargs.pop('header_font_size',self.header_font_color)
        fill_index = kwargs.pop('fill_index',self.fill_index)
        bold_index = kwargs.pop('bold_index',self.bold_index)
        index_font_size = kwargs.pop('index_font_size',self.index_font_size)
        index_font_color = kwargs.pop('index_font_color',self.index_font_color)
        fill_color = kwargs.pop('fill_color',self.fill_color)
        row_banding = kwargs.pop('row_banding',self.row_banding)
        column_banding = kwargs.pop('column_banding',self.column_banding)
        first_row = kwargs.pop('first_row',self.first_row)
        first_col = kwargs.pop('first_col',self.first_col)
        last_row = kwargs.pop('last_row',self.last_row)
        last_col = kwargs.pop('last_col',self.last_col)
        merge_indices = kwargs.pop('merge_indices',self.merge_indices)
        center_merge = kwargs.pop('center_merge',self.center_merge)
        row_height = kwargs.pop('row_height',self.row_height)
        if (fill_header or bold_header or not header_font_size == None or not header_font_color == None) and not self.header:
            raise ValueError("Cannot style DataFrame header when table.header attribute is False. \
            First set `table.header = True` to include the DataFrame header in the output table.")
        elif self.header:
            self.style_index(axis=1,
                             fill=fill_header,
                             bold=bold_header,
                             font_size=header_font_size,
                             font_color=header_font_color,
                             fill_color=fill_color,
                             merge_indices=merge_indices,
                             center_merge=center_merge)
        if (fill_index or bold_index or not index_font_size == None or not index_font_color == None) and not self.index:
            raise ValueError("Cannot style DataFrame index when table.index attribute is False. \
            First set `table.header = True` to include the DataFrame header in the output table.")
        elif self.index:
            self.style_index(axis=0,
                             fill=fill_index,
                             bold=bold_index,
                             font_size=index_font_size,
                             font_color=index_font_color,
                             fill_color=fill_color,
                             merge_indices=merge_indices,
                             center_merge=center_merge)
        table = pptFunctions.set_row_height(table,
                                         row_height=row_height)
        table.horz_banding = row_banding
        table.vert_banding = column_banding
        table.first_row = first_row
        table.first_col = first_col
        table.last_row = last_row
        table.last_col = last_col

    def convert(self, **kwargs):
        """Perform the conversion from DataFrame values to table cells.

        Returns
        -------
        data: pandas.DataFrame
            Formatted and transformed DataFrame exported to PowerPoint.

        Notes
        -----
        This method calls the individual processing methods in sequence, \
        then builds the PowerPoint table by inserting values cell by cell.
        """
        self.insert_table()
        table = self.shape.table
        data = self.df.copy()
        data = self.format_index(data=data, axis=0)
        data = self.format_index(data=data, axis=1)
        if self.column_totals:
            data = self.add_totals(data=data, axis=0,
            totals_label=self.column_totals_label,
            totals_method=self.column_totals_method,
            totals_aggmap=self.column_totals_aggmap)
        if self.row_totals:
            data = self.add_totals(data=data, axis=1,
            totals_label=self.row_totals_label,
            totals_method=self.row_totals_method,
            totals_aggmap=self.row_totals_aggmap)
        data = self.format_values(data=data)
        data = self.transform(data=data)
        data = data.fillna(self.na_rep)
        for (row,col),x in np.ndenumerate(data.values):
            c = table.cell(row,col)
            c.text = x if isinstance(x,str) else str(x)
            c = pptFunctions.format_cell(c,
                                      font_size=self.font_size,
                                      font_color=self.font_color,
                                      font_name=self.font_name,
                                      cell_margins=self.cell_margins)
        self.style_table()
        self.align_columns()
        return data
